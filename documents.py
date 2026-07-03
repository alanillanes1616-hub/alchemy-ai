"""
documents.py — Extracción de texto de múltiples formatos de archivo.

Cada extractor devuelve una lista de dicts:
    [{"page": int, "text": str, "source": str}, ...]
"""

import io
import csv
from pathlib import Path
from typing import Protocol

import pandas as pd
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from utils import clean_text, extension, get_logger

logger = get_logger("documents")


# ── Protocolo base ────────────────────────────────────────────────────────────

class Extractor(Protocol):
    def extract(self, content: bytes, filename: str) -> list[dict]:
        ...


# ── Extractores por formato ───────────────────────────────────────────────────

class PDFExtractor:
    def extract(self, content: bytes, filename: str) -> list[dict]:
        pages = []
        try:
            reader = PdfReader(io.BytesIO(content))
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                text = clean_text(text)
                if text:
                    pages.append({"page": i + 1, "text": text, "source": filename})
        except Exception as exc:
            logger.error("PDF extraction failed for %s: %s", filename, exc)
        return pages


class DOCXExtractor:
    def extract(self, content: bytes, filename: str) -> list[dict]:
        pages = []
        try:
            doc = DocxDocument(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Agrupa en bloques de ~50 párrafos simulando "páginas"
            block_size = 50
            for i in range(0, len(paragraphs), block_size):
                chunk = "\n".join(paragraphs[i : i + block_size])
                text = clean_text(chunk)
                if text:
                    pages.append({
                        "page": i // block_size + 1,
                        "text": text,
                        "source": filename,
                    })
        except Exception as exc:
            logger.error("DOCX extraction failed for %s: %s", filename, exc)
        return pages


class TXTExtractor:
    def extract(self, content: bytes, filename: str) -> list[dict]:
        try:
            text = clean_text(content.decode("utf-8", errors="replace"))
            if text:
                return [{"page": 1, "text": text, "source": filename}]
        except Exception as exc:
            logger.error("TXT extraction failed for %s: %s", filename, exc)
        return []


class MarkdownExtractor(TXTExtractor):
    """Markdown es texto plano; reutiliza TXTExtractor."""
    pass


class PPTXExtractor:
    def extract(self, content: bytes, filename: str) -> list[dict]:
        pages = []
        try:
            prs = Presentation(io.BytesIO(content))
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                text = clean_text("\n".join(texts))
                if text:
                    pages.append({"page": i + 1, "text": text, "source": filename})
        except Exception as exc:
            logger.error("PPTX extraction failed for %s: %s", filename, exc)
        return pages


class XLSXExtractor:
    def extract(self, content: bytes, filename: str) -> list[dict]:
        pages = []
        try:
            xls = pd.ExcelFile(io.BytesIO(content))
            for i, sheet in enumerate(xls.sheet_names):
                df = xls.parse(sheet)
                text = clean_text(df.to_string(index=False))
                if text:
                    pages.append({"page": i + 1, "text": text, "source": filename})
        except Exception as exc:
            logger.error("XLSX extraction failed for %s: %s", filename, exc)
        return pages


class CSVExtractor:
    def extract(self, content: bytes, filename: str) -> list[dict]:
        try:
            text = clean_text(content.decode("utf-8", errors="replace"))
            if text:
                return [{"page": 1, "text": text, "source": filename}]
        except Exception as exc:
            logger.error("CSV extraction failed for %s: %s", filename, exc)
        return []


# ── Registro de extractores ───────────────────────────────────────────────────

_EXTRACTORS: dict[str, Extractor] = {
    "pdf":  PDFExtractor(),
    "docx": DOCXExtractor(),
    "txt":  TXTExtractor(),
    "md":   MarkdownExtractor(),
    "pptx": PPTXExtractor(),
    "xlsx": XLSXExtractor(),
    "csv":  CSVExtractor(),
}


# ── Función pública ───────────────────────────────────────────────────────────

def extract_document(content: bytes, filename: str) -> list[dict]:
    """
    Extrae texto de un archivo dado su contenido en bytes.

    Returns:
        Lista de dicts con keys: page, text, source.
        Lista vacía si el formato no está soportado o falla la extracción.
    """
    ext = extension(filename)
    extractor = _EXTRACTORS.get(ext)
    if not extractor:
        logger.warning("Unsupported extension: %s (%s)", ext, filename)
        return []

    pages = extractor.extract(content, filename)
    logger.info("Extracted %d pages from %s", len(pages), filename)
    return pages
